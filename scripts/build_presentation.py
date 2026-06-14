#!/usr/bin/env python3
"""Build .pptx presentation from docs/presentation.md"""

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
INPUT = ROOT / "docs" / "presentation.md"
OUTPUT_DIR = ROOT / "docs" / "dist"

# Brand colors
INDIGO = RGBColor(0x43, 0x38, 0xCA)
DARK = RGBColor(0x1E, 0x1B, 0x4B)
GRAY = RGBColor(0x37, 0x41, 0x51)
LIGHT_GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF3, 0xFF)
TABLE_HEADER_BG = RGBColor(0x43, 0x38, 0xCA)
TABLE_ALT_BG = RGBColor(0xF9, 0xFA, 0xFB)


def parse_slides(content: str) -> list[str]:
    """Split markdown into slides on --- separators."""
    if content.startswith("---"):
        parts = content.split("---", 2)
        if len(parts) >= 3:
            content = parts[2]

    slides = re.split(r"\n\s*---\s*\n", content.strip())
    return [s.strip() for s in slides if s.strip()]


def parse_slide_content(md: str) -> dict:
    """Parse a slide's markdown into structured content."""
    lines = md.split("\n")
    result = {
        "title": "",
        "subtitle": "",
        "elements": [],  # list of (type, content)
    }

    in_code = False
    code_lines = []
    in_table = False
    table_rows = []

    for line in lines:
        # Code blocks
        if line.startswith("```"):
            if in_code:
                result["elements"].append(("code", "\n".join(code_lines)))
                code_lines = []
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_lines.append(line)
            continue

        # Tables
        if "|" in line and not line.startswith(">"):
            cells = [c.strip() for c in line.split("|")[1:-1]]
            if all(re.match(r"^[-:]+$", c) for c in cells):
                continue
            if not in_table:
                in_table = True
            table_rows.append(cells)
            continue
        elif in_table:
            result["elements"].append(("table", table_rows))
            table_rows = []
            in_table = False

        # Headers
        if line.startswith("# "):
            result["title"] = strip_md(line[2:])
            continue
        if line.startswith("## "):
            result["subtitle"] = strip_md(line[3:])
            continue
        if line.startswith("### "):
            result["elements"].append(("subheading", strip_md(line[4:])))
            continue

        # Blockquotes
        if line.startswith("> "):
            result["elements"].append(("quote", strip_md(line[2:])))
            continue

        # List items
        if line.startswith("- "):
            result["elements"].append(("bullet", strip_md(line[2:])))
            continue
        if re.match(r"^\d+\. ", line):
            text = re.sub(r"^\d+\. ", "", line)
            result["elements"].append(("bullet", strip_md(text)))
            continue

        # Regular text
        if line.strip():
            result["elements"].append(("text", strip_md(line)))

    if in_table and table_rows:
        result["elements"].append(("table", table_rows))
    if in_code and code_lines:
        result["elements"].append(("code", "\n".join(code_lines)))

    return result


def strip_md(text: str) -> str:
    """Remove markdown formatting, keeping text."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", text)
    return text


def has_bold(text: str) -> list[tuple[str, bool]]:
    """Split text into segments with bold flag."""
    segments = []
    parts = re.split(r"(\*\*.+?\*\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            clean = part[2:-2]
            clean = re.sub(r"`(.+?)`", r"\1", clean)
            segments.append((clean, True))
        else:
            clean = re.sub(r"\*(.+?)\*", r"\1", part)
            clean = re.sub(r"`(.+?)`", r"\1", clean)
            clean = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", clean)
            segments.append((clean, False))
    return segments


def add_gradient_bg(slide):
    """Add a subtle gradient background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = WHITE
    fill.gradient_stops[1].color.rgb = LIGHT_BG


def add_title_slide(prs, content):
    """Create a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_gradient_bg(slide)

    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content["title"]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = INDIGO
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    if content["subtitle"]:
        top = Inches(4.0)
        txBox = slide.shapes.add_textbox(left, top, width, Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content["subtitle"]
        p.font.size = Pt(22)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.LEFT

    # Body elements
    body_top = Inches(4.8) if content["subtitle"] else Inches(4.0)
    add_body_elements(slide, content["elements"], body_top)


def add_content_slide(prs, content):
    """Create a standard content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_gradient_bg(slide)

    # Title
    left = Inches(0.8)
    top = Inches(0.4)
    width = Inches(8.4)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content["title"]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INDIGO

    # Subtitle
    body_top = Inches(1.3)
    if content["subtitle"]:
        txBox = slide.shapes.add_textbox(left, body_top, width, Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content["subtitle"]
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        body_top = Inches(2.0)

    add_body_elements(slide, content["elements"], body_top)


def add_body_elements(slide, elements, start_top):
    """Add body content elements to a slide."""
    left = Inches(0.8)
    width = Inches(8.4)
    current_top = start_top

    for etype, econtent in elements:
        if current_top > Inches(7.0):
            break

        if etype == "subheading":
            txBox = slide.shapes.add_textbox(left, current_top, width, Inches(0.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = econtent
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = DARK
            current_top += Inches(0.45)

        elif etype == "text":
            txBox = slide.shapes.add_textbox(left, current_top, width, Inches(0.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            segments = has_bold(econtent)
            p.clear()
            for text, bold in segments:
                run = p.add_run()
                run.text = text
                run.font.size = Pt(14)
                run.font.bold = bold
                run.font.color.rgb = GRAY
            current_top += Inches(0.4)

        elif etype == "bullet":
            txBox = slide.shapes.add_textbox(
                left + Inches(0.3), current_top, width - Inches(0.3), Inches(0.45)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            segments = has_bold(econtent)
            p.clear()
            bullet_run = p.add_run()
            bullet_run.text = "•  "
            bullet_run.font.size = Pt(13)
            bullet_run.font.color.rgb = INDIGO
            for text, bold in segments:
                run = p.add_run()
                run.text = text
                run.font.size = Pt(13)
                run.font.bold = bold
                run.font.color.rgb = GRAY
            current_top += Inches(0.35)

        elif etype == "quote":
            # Add accent bar
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left + Inches(0.2),
                current_top + Inches(0.05),
                Inches(0.08),
                Inches(0.4),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = INDIGO
            shape.line.fill.background()

            txBox = slide.shapes.add_textbox(
                left + Inches(0.5), current_top, width - Inches(0.5), Inches(0.5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = econtent
            p.font.size = Pt(12)
            p.font.italic = True
            p.font.color.rgb = INDIGO
            current_top += Inches(0.55)

        elif etype == "code":
            lines = econtent.split("\n")
            box_height = min(Inches(0.25 * len(lines) + 0.3), Inches(2.5))
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left,
                current_top,
                width,
                box_height,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = DARK
            shape.line.fill.background()

            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.2)
            tf.margin_top = Inches(0.15)
            display_lines = lines[:10]
            if len(lines) > 10:
                display_lines.append("...")
            for i, code_line in enumerate(display_lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = code_line
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(0xE0, 0xE7, 0xFF)
            current_top += box_height + Inches(0.2)

        elif etype == "table":
            rows_data = econtent
            if not rows_data:
                continue
            n_rows = min(len(rows_data), 10)
            n_cols = len(rows_data[0])

            table_height = Inches(0.35 * n_rows + 0.1)
            if current_top + table_height > Inches(7.2):
                table_height = Inches(7.2) - current_top
                n_rows = min(n_rows, int((table_height / Inches(0.35))))

            table_shape = slide.shapes.add_table(
                n_rows, n_cols, left, current_top, width, table_height
            )
            table = table_shape.table

            for col_idx in range(n_cols):
                table.columns[col_idx].width = int(width / n_cols)

            for row_idx in range(n_rows):
                for col_idx in range(n_cols):
                    cell = table.cell(row_idx, col_idx)
                    cell_text = rows_data[row_idx][col_idx] if col_idx < len(rows_data[row_idx]) else ""
                    cell.text = strip_md(cell_text)

                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(10)
                        if row_idx == 0:
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = WHITE
                        else:
                            paragraph.font.color.rgb = GRAY

                    if row_idx == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = TABLE_HEADER_BG
                    elif row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = TABLE_ALT_BG

            current_top += table_height + Inches(0.2)


def build_pptx(slides_md: list[str]) -> Presentation:
    """Build a PowerPoint presentation from parsed slide content."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for i, slide_md in enumerate(slides_md):
        content = parse_slide_content(slide_md)
        if i == 0:
            add_title_slide(prs, content)
        else:
            add_content_slide(prs, content)

    return prs


def main():
    if not INPUT.exists():
        print(f"Error: {INPUT} not found", file=sys.stderr)
        sys.exit(1)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    content = INPUT.read_text()
    slides = parse_slides(content)
    prs = build_pptx(slides)

    output_file = OUTPUT_DIR / "SynthFocus-Pitch.pptx"
    prs.save(str(output_file))
    print(f"Built presentation: {output_file}")
    print(f"  Slides: {len(slides)}")
    print(f"  Size: {output_file.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
